import configparser
import contextlib
from ctypes import windll
from io import BytesIO
import itertools
from json import dump
from multiprocessing import Pool, freeze_support
import os
import re
from signal import signal, SIGINT, SIG_IGN
from subprocess import run, DEVNULL
import sys
import time
import webbrowser

from colorama import init, Fore, Style
import cursor
import cv2
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE  # type: ignore
from pptx.slide import Slide
from pptx.util import Inches, lazyproperty  # TODO: Remove when done inspecting
from pywintypes import com_error
import requests
import win32com.client

from client import ProgramStatus, fetch_latest_version
from client import download_avatar
from exceptions import Error, ErrorType, print_exception_hook
from utils import is_number, as_int, hex_to_rgb, parse_version
from utils import app_dir, abs_path
from utils import inp, console_style, ProgressBar
from vba.macros import module1_bas


def replace_text(slide: Slide, df, i, avatar_mode) -> Slide:
    """Replaces and formats text."""
    cols = df.columns.values.tolist() + ['p']

    for shape in slide.shapes:  # type: ignore
        if not shape.has_text_frame or '{' not in shape.text:
            continue

        text_frame = shape.text_frame

        for run in itertools.chain.from_iterable(
            [p.runs for p in text_frame.paragraphs]):

            for search_str in (set(re.findall(r'(?<={)(.*?)(?=})', run.text))
                               .intersection(cols)):

                # Profile picture
                if search_str == 'p':
                    # Test cases
                    if 'uid' not in cols or not avatar_mode:
                        continue

                    if pd.isnull(df['uid'].iloc[i]):
                        run.text = ''
                        continue

                    # Extract effect index and remove {p}
                    effect = as_int(run.text.strip()[3:])
                    run.text = ''

                    uid = str(df['uid'].iloc[i]).strip().replace('_', '')

                    og_path = avatar_dir + '_' + uid + '.png'
                    img_path = avatar_dir + str(effect) + '_' + uid + '.png'

                    if not os.path.isfile(og_path):
                        continue

                    if is_number(effect):
                        img = cv2.imread(og_path)
                        match effect:  # TODO: Add more effects in the future
                            case 1:
                                img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

                        cv2.imwrite(img_path, img)

                    new_shape = slide.shapes._fget().add_picture(  # TODO
                        img_path, shape.left, shape.top,
                        shape.width, shape.height
                    )

                    new_shape.auto_shape_type = MSO_SHAPE.OVAL
                    old = shape._element
                    new = new_shape._element
                    old.addnext(new)
                    old.getparent().remove(old)
                    continue

                # Actual text
                repl = str(df[search_str].iloc[i])
                repl = repl if repl != 'nan' else ''  # Replace nan with empty

                run_text = run.text

                if search_str.startswith(starts):
                    run.text = repl
                else:
                    run.text = run.text.replace('{' + search_str + '}', repl)

                # Replace image links
                pattern = r'\<\<(.*?)\>\>'  # Look for <<image_links.{ext}>>
                img_link = re.findall(pattern, run.text)

                if len(img_link) > 0:
                    try:
                        img = BytesIO(requests.get(img_link[0]).content)
                        pil = Image.open(img)

                        im_width = shape.height / pil.height * pil.width
                        new_shape = slide.shapes.add_picture(  # type: ignore
                            img, shape.left + (shape.width - im_width)/2,
                            shape.top, im_width, shape.height
                        )

                        old = shape._element.addnext(new_shape._element)

                        run.text = re.sub(pattern, '', run.text)
                        text_frame.margin_left = Inches(5.2)
                    except Exception:
                        Error(
                            'Could not load the following image from '
                           f'slide {i + 1}, sheet {df["sheet"].iloc[0]}.',
                           f'{img_link[0]}',
                            'Please check your internet connection and verify '
                            'that the link directs to an image file, which '
                            'usually ends in an image extension like .png.',
                            err_type=ErrorType.WARNING).throw()

                # Color formatting
                if not search_str.startswith(starts):
                    continue

                # Check RGB
                if (run.font.color.type == MSO_COLOR_TYPE.RGB and  # type: ignore
                    run.font.color.rgb not in [
                        RGBColor(0, 0, 0), RGBColor(255, 255, 255)]):
                    continue

                for ind, val in enumerate(range_list):
                    if is_number(repl) and float(repl) >= val:
                        if run_text.endswith('1'):
                            run.font.color.rgb = RGBColor(*scheme_alt[ind])
                        else:
                            run.font.color.rgb = RGBColor(*scheme[ind])
                        break
    return slide


if __name__ == '__main__':
    version_tag = '3.0'

    # Section A: Fix console issues
    freeze_support()          # Multiprocessing freeze support
    signal(SIGINT, SIG_IGN)   # Handle KeyboardInterrupt

    # Disable QuickEdit and Insert mode
    kernel32 = windll.kernel32
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x00|0x100))

    sys.excepthook = print_exception_hook  # Avoid exiting program on exception
    init()                                 # Enable ANSI escape sequences
    cursor.hide()                          # Hide cursor


    # Section B: Get current directories and files
    output_dir = abs_path('output')
    avatar_dir = abs_path('avatars')

    if missing := [f for f in (
            'settings.ini',
            'data.xlsx',
            'template.pptm',
            'token.txt',
        ) if not os.path.exists(abs_path(f))]:
        Error(40).throw(app_dir, '\n'.join(missing))


    # Section C: Load user configurations
    config = configparser.ConfigParser()
    config.read(abs_path('settings.ini'))

    # Store config variables as local variables
    range_list = config['FORMATTING']['ranges'][::-1]
    scheme = config['FORMATTING']['scheme'][::-1]
    scheme_alt = config['FORMATTING']['scheme_alt'][::-1]
    starts = config['FORMATTING']['trigger_word']
    avatar_mode = config['AVATARS']['avatar_mode']
    last_clear = config['AVATARS']['last_clear_avatar_cache']

    with open(abs_path('token.txt')) as f:
        token_list = f.read().splitlines()
        token_list = [i.strip() for i in token_list if len(i) > 62]

    if not token_list and avatar_mode:
        Error(21).throw()

    scheme = list(map(hex_to_rgb, scheme))
    scheme_alt = list(map(hex_to_rgb, scheme_alt))


    # Section D: Check for updates
    status = None
    repo_url = 'https://github.com/banz04/mic-drop-results/'

    with contextlib.suppress(requests.exceptions.ConnectionError, KeyError):
        if config['PROGRAM']['update_check']:
            # Fetch the latest version and the summary of the update
            latest_tag, summary = fetch_latest_version()

            latest, current = parse_version(
                latest_tag, version_tag
            )

            if latest > current:
                status = ProgramStatus.UPDATE_AVAILABLE

                console_style(Fore.YELLOW, Style.BRIGHT)
                print(f'Update v{latest_tag}')

                console_style(Style.NORMAL)
                print(summary)

                release_url = f'{repo_url}releases/latest/'
                print(release_url + '\n')
                console_style()

                webbrowser.open(release_url, new=2)

            elif latest < current:
                status = ProgramStatus.BETA
            else:
                status = ProgramStatus.UP_TO_DATE

    # Print a header containing information about the program

    # Normal:       Mic Drop Results (vX.1) [latest]
    #               https://github.com/banz04/mic-drop-results


    # With update:  Update vX.1
    #               A summary of the update will appear in this line.
    #               https://github.com/banz04/mic-drop-results/releases/latest/
    #
    #               Mic Drop Results (vX.0) [update available]

    status_msg = f' [{status.value}]' if status else ''
    print(f'Mic Drop Results (v{version_tag}){status_msg}')
    console_style()

    if status != ProgramStatus.UPDATE_AVAILABLE:
    # When an update is available, the download link is already shown above.
    # To avoid confusion, we only print one link at a time.
        print(repo_url)


    # Section E: Process the data
    xls = pd.ExcelFile('data.xlsx')

    sheetnames_raw = xls.sheet_names
    sheetnames = [re.sub(r'[\\\/:"*?<>|]+',  # Forbidden file name characters
        '', sheet) for sheet in sheetnames_raw]

    data = {}

    db_list = []
    for sheet in sheetnames_raw:
        if sheet.startswith('_'):
            df = pd.read_excel(xls, sheet)

            # Validate shape
            if df.empty or df.shape < (1, 2):
                continue

            db_list.append(df)

    SHARING_VIOLATION = '\033[33mNOTE: Please exit the program before modifying data.xlsx or ' \
                        'Microsoft Excel will throw a Sharing Violation error.\033[39m'

    for i, sheet in enumerate(sheetnames_raw):
        df = pd.read_excel(xls, sheet)

        # Validate shape
        if df.empty or df.shape < (1, 2):
            continue

        # Exclude database sheets
        if sheet.startswith('_'):
            continue

        # Exclude sheets with first two columns where data is not numeric
        if sum(df.iloc[:, i].dtype.kind in 'biufc' for i in range(2)) < 2:
            Error(f'Invalid data type. The following rows of {sheet} contain strings '
                   'instead of the supposed numeric data type within the first two columns. '
                   'The sheet will be excluded if you proceed on.',

                df[~df.iloc[:, :2].applymap(np.isreal).all(1)],

                err_type=ErrorType.WARNING
            ).throw()

            continue

        # Replace NaN values within the first two columns with 0
        if df.iloc[:, :2].isnull().values.any():
            Error(f'The following rows of {sheet} contain empty values '
                'within the first two columns.',

                df[df.iloc[:, :2].isnull().any(axis=1)],

                'You may exit this program and modify your data or proceed on with '
                'these empty values substituted with 0.', SHARING_VIOLATION,

                err_type=ErrorType.WARNING
            ).throw()

            df.iloc[:, :2] = df.iloc[:, :2].fillna(0)

        # Check for cases where avg and std are the same (hold the same rank)
        df['r'] = pd.DataFrame(zip(df.iloc[:, 0], df.iloc[:, 1] * -1)) \
                    .apply(tuple, axis=1).rank(method='min', ascending=False).astype(int)

        # Sort the slides
        df = df.sort_values(by='r', ascending=True)

        # Remove .0 from whole numbers
        format_number = lambda x: str(int(x)) if x % 1 == 0 else str(x)
        df.loc[:, df.dtypes == float] = df.loc[:, df.dtypes == float].applymap(format_number)

        # Replace {sheet} with sheet name
        df['sheet'] = sheet

        # Merge contestant database
        clean_name = lambda x: x.str.lower().str.strip() if(x.dtype.kind == 'O') else x
        if db_list:
            for db in db_list:
                df_cols = df.columns.values.tolist()
                db_cols = db.columns.values.tolist()
                merge_col = db_cols[0]

                # Use merge for non-existing columns
                df = df.merge(db[[merge_col] + [i for i in db_cols if i not in df_cols]],
                    left_on=clean_name(df[merge_col]), right_on=clean_name(db[merge_col]), how='left')

                df.loc[:, merge_col] = df[merge_col + '_x']
                df.drop(['key_0', merge_col + '_x', merge_col + '_y'], axis=1, inplace=True)

                # Use update for existing columns
                df['update_index'] = clean_name(df[merge_col])
                df = df.set_index('update_index')

                db['update_index'] = clean_name(db[merge_col])
                db = db.set_index('update_index')
                db_cols.remove(merge_col)

                update_cols = [i for i in db_cols if i in df_cols]
                df.update(db[update_cols])
                df.reset_index(drop=True, inplace=True)

        if 'uid' not in df.columns.values.tolist(): avatar_mode = 0

        # Fill in missing templates
        df['template'].fillna(1, inplace=True)

        data[sheetnames[i]] = df

    if not data:
        Error(f'No valid sheet found in {app_dir}data.xlsx').throw()


    # Section F: Generate PowerPoint slides
    print('\nGenerating slides...')
    print('Please do not click on any PowerPoint windows that may show up in the process.\n')

    # Kill all PowerPoint instances
    run('TASKKILL /F /IM powerpnt.exe', stdout=DEVNULL, stderr=DEVNULL)

    # Open template presentation
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(avatar_dir, exist_ok=True)

    # Clear cache
    if time.time() - last_clear > 1800:  # Clears every hour
        for f in os.scandir(avatar_dir):
            os.unlink(f)

        # Update last clear time
        with open('config.json', 'w') as f:
            config['last_clear_avatar_cache'] = int(time.time())
            dump(config, f, indent=4)

    # Download avatars with parallel processing
    attempt = 0
    pool = Pool(3)

    while avatar_mode:
        uid_list = []

        for df in data.values():
            if df['uid'].dtype.kind in 'biufc':
                Error('The \'uid\' column has a numeric data type instead of the supposed string data type.',
                      'Please exit the program and add an underscore before every user ID.', SHARING_VIOLATION).throw()

            uid_list += [id for id in df['uid'] if not pd.isnull(id) and not os.path.isfile(avatar_dir + id.strip() + '.png')]

        if len(uid_list) == 0:
            break

        if attempt > 0 and attempt <= 3:
            print(f'Unable to download the profile pictures of the following users. Retrying {attempt}/3',
                    uid_list, sep='\n', end='\n\n')
        elif attempt > 3:
            Error(23).throw(str(uid_list), err_type=ErrorType.WARNING)

        pool.starmap(download_avatar, zip(uid_list,
            [avatar_dir] * len(uid_list), itertools.islice(itertools.cycle(token_list), len(uid_list))))

        attempt += 1

    pool.close()
    pool.join()

    for k, df in data.items():
        bar = ProgressBar(8, title=k, max_title_length=max(map(len, data.keys())))

        # Open template presentation
        bar.set_description('Opening template.pptm')
        ppt = win32com.client.Dispatch('PowerPoint.Application')
        ppt.Presentations.Open(f'{app_dir}template.pptm')
        bar.add()

        # Import macros
        bar.set_description('Importing macros')

        try:
            ppt.VBE.ActiveVBProject.VBComponents.Import(module1_bas)
        except com_error as e:
            if e.hresult == -2147352567:  # type: ignore
            # Trust access settings not yet enabled
                Error(41).throw()
            else:
                raise e

        bar.add()

        # Duplicate slides
        bar.set_description('Duplicating slides')
        slides_count = ppt.Run('Count')

        # Duplicate slides
        for t in df.loc[:, 'template']:
            if as_int(t) not in range(1, slides_count + 1):
                Error(f'Template {t} does not exist (error originated from the following sheet: {k}).',
                      f'Please exit the program and modify the \'template\' column of {k}.', SHARING_VIOLATION).throw()

            ppt.Run('Duplicate', t)

        bar.add()

        # Delete template slides when done
        ppt.Run('DelSlide', *range(1, slides_count + 1))
        bar.add()

        # Save as output file
        bar.set_description('Saving templates')
        output_filename = f'{k}.pptx'

        ppt.Run('SaveAs', f'{output_dir}{output_filename}')
        bar.add()

        run('TASKKILL /F /IM powerpnt.exe', stdout=DEVNULL, stderr=DEVNULL)
        bar.add()

        # Replace text
        bar.set_description('Filling in judging data')
        prs = Presentation(output_dir + output_filename)

        for i, slide in enumerate(prs.slides):
            replace_text(slide, df, i, avatar_mode)
        bar.add()

        # Save
        bar.set_description(f'Saving as {output_dir + output_filename}')
        prs.save(output_dir + output_filename)
        bar.add()


    # Section G: Launch the file
    print(f'\nExported to {output_dir}')

    # Enable QuickEdit
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x40|0x100))

    inp('Press Enter to open the output folder...')
    os.startfile(output_dir)
