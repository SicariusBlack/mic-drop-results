# Copyright 2023 Phan Huy

# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.

import atexit
from concurrent.futures import ThreadPoolExecutor
import contextlib
from io import BytesIO
import itertools
from multiprocessing import freeze_support
import os
from signal import signal, SIGINT, SIG_IGN
from subprocess import check_call, run, DEVNULL
import sys
import threading
import time
import warnings
import webbrowser

import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.slide import Slide
from pptx.util import Cm
from pywintypes import com_error
from rich.padding import Padding
import requests
import win32com.client

from client import (
    ProgramStatus,
    fetch_avatar,
    fetch_latest_version,
    download_avatars,
    _get_download_banner,
)
from compiled_regex import *
from config import Config
import constants
from constants import *
from errors import Error, ErrorType, print_exception_hook
from exceptions import *
from utils import (
    inp,
    enable_console,
    disable_console,
    chunk,
    is_number,
    as_type,
    hex_to_rgb,
    parse_version,
    parse_coef,
    clean_name,
    abs_dir,
    get_avatar_dir,
    artistic_effect,
)
from vba.macros import module1_bas


def _replace_avatar(slide: Slide, shape, run, *, uid: str) -> None:
    """Replaces avatar element on slide with avatar.

    Args:
        slide (Slide): _description_
        shape: _description_
        run: _description_
        uid (str): _description_
    """
    effect_id = parse_coef(
        run.text, field_name="p"
    )  # parse effect id from the field's coefficient
    run.text = ""  # reset text box to empty

    avatar_og_dir = get_avatar_dir(uid)  # get avatar without effect directory
    if not avatar_og_dir.is_file():
        return

    # Add avatar to slide
    avatar_dir = artistic_effect(avatar_og_dir, effect=effect_id)
    new_shape = slide.shapes.add_picture(  # type: ignore
        str(avatar_dir), shape.left, shape.top, shape.width, shape.height
    )
    new_shape.auto_shape_type = MSO_SHAPE.OVAL
    old = shape._element
    new = new_shape._element
    old.addnext(new)
    old.getparent().remove(old)


def _replace_text(run, field_name, *, text: str) -> None:
    if text == "nan":  # change numpy's nan back to empty value
        text = ""

    if field_name.startswith(cfg.trigger_word):  # apply to {triggerword_blahblah}
        # Find breakpoints and apply conditional formatting to numbers
        for ind, breakpoint in enumerate(cfg.ranges[::-1]):
            if not is_number(text):
                break

            if float(text) >= breakpoint:
                if parse_coef(run.text, field_name=field_name) == 0:
                    run.font.color.rgb = RGBColor(*scheme[::-1][ind])
                else:
                    run.font.color.rgb = RGBColor(*scheme_alt[::-1][ind])
                break

        run.text = text
    else:
        # Replace text normally
        run.text = run.text.replace("{" + field_name + "}", text)


def _replace_image_url(shape, p, run) -> None:
    if img_url := match_url.findall(run.text):  # find image urls
        with contextlib.suppress(Exception):
            margin_left = _insert_image(shape, img_url=img_url[0])
            run.text = run.text.replace(img_url[0], "")
            # With some measurements we can obtain 12.47 cm = 4490850
            # 1 cm = 360132.3175621492
            shape.text_frame.margin_left = Cm(margin_left / 360132.3175621492)
            p.alignment = PP_ALIGN.LEFT


def _insert_image(shape, *, img_url: str) -> float:
    """Inserts an image on top of a shape on slide.

    Args:
        shape: the shape to fit the image.
        img_url (str): the URL of the image on the internet.

    Returns:
        float: the left margin to indent the remaining text.
    """
    im_bytes = BytesIO(requests.get(img_url).content)
    img = Image.open(im_bytes)

    height = shape.height
    width = height / img.height * img.width
    left = shape.left + (shape.width - width) / 2
    top = shape.top

    new_shape = slide.shapes.add_picture(  # type: ignore
        im_bytes, left, top, width, height
    )
    shape._element.addnext(new_shape._element)

    return left + width  # left margin for the remaining text


def fill_slide(slide: Slide, data: dict[str, str]) -> None:
    columns = set([*data] + ["p"])

    for shape in slide.shapes:  # type: ignore
        if not shape.has_text_frame:
            continue

        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                for field_name in match_field_name.findall(run.text):
                    field_name = field_name.lstrip("__")
                    if field_name not in columns:
                        continue

                    # Replace {p} with avatar
                    if field_name == "p":
                        _replace_avatar(slide, shape, run, uid=data["uid"])
                        break

                    # Replace text
                    _replace_text(run, field_name, text=data[field_name])

                    _replace_image_url(shape, p, run)


def preview_df(
    df: pd.DataFrame,
    filter_series: pd.Series | None = None,
    *,
    n_cols: int,
    n_cols_ext: int = 5,
    highlight: bool = True,
    words_to_highlight: list[str | None] | None = None,
) -> str:
    """Formats and returns a string preview of the dataframe.

    Args:
        df: the dataframe to format.
        filter_series (optional): the boolean series used to select
            rows. Pass None to include all rows in the preview. Defaults
            to None.
        n_cols: the number columns to format (starts at first column).
            If n_cols is less than the number of columns in the dataframe,
            the preview will be a snippet, which shows ellipses at the
            end of every line).
        n_cols_ext (optional): the number of extra columns to preview
            alongside with the formatted columns. Defaults to 5.
        highlight (optional): enable value highlighting. Defaults to
            True.
        words_to_highlight (optional): list of words
            to highlight in red (only effective to values in the first
            n_cols columns). List None if the value to highlight is
            numpy.nan. Pass None to highlight nothing. Defaults to None.

    Returns:
        str: the formatted dataframe as a string.
    """
    if words_to_highlight is None:
        words_to_highlight = []

    df = df.copy(deep=True)
    if filter_series is not None:
        df = df[filter_series]
    df.index += 2  # reflect row index as displayed in Excel

    # Only show the first few columns for preview
    df = df.iloc[:, : min(n_cols + n_cols_ext, len(df.columns))]

    # Replace text_to_highlight with ⁅text_to_highlight⁆
    prefix, suffix = "⁅", "⁆"  # arbitrary symbols, must be single length
    highlight_str = lambda x: f"{prefix}{x}{suffix}"
    for word in words_to_highlight:
        if word is None:
            df.iloc[:, :n_cols] = df.iloc[:, :n_cols].fillna(highlight_str("NaN"))
        else:
            df.iloc[:, :n_cols] = df.iloc[:, :n_cols].replace(word, highlight_str(word))

    preview = repr(df.head(8)) if n_cols < len(df.columns) else repr(df)

    # Highlight text_to_highlight
    preview = preview.replace(prefix, "[red]  ").replace(suffix, "[/red]")

    # Highlight column names
    for n in range(n_cols):
        col = df.columns.tolist()[n]
        preview = preview.replace(" " + col, f" [red]{col}[/red]", 1)

    # Add ... at the end of each line if preview is a snippet
    if n_cols < len(df.columns):
        preview = preview.replace("\n", "  ...\n") + "  ..."

    # Bold first row
    preview = "[b]" + preview.replace("\n", "[/b]\n", 1)

    if not highlight:
        preview = preview.replace("[red]", "")
    return preview


def _import_avatars():
    failed = False  # whether the download task has failed
    max_attempt = 5  # maximum number of attempts
    has_task = False  # whether a download task exists (to skip avatar download banner)
    uids_unknown = []  # uids that failed to download

    for attempt in range(1, max_attempt + 1):
        uids = []  # uids that are not downloaded yet
        for df in groups.values():
            if df["__uid"].dtype.kind in "biufc":  # if uid column is numeric
                Error(70).throw()

            for id in df["__uid"]:
                if not (
                    pd.isnull(id)  # skip nan values
                    or get_avatar_dir(id).is_file()  # skip if already downloaded
                    or id in uids  # skip if uid already in queue
                    or id in uids_unknown  # skip if already in the unknown list
                ):
                    uids.append(id)
        if not uids:  # if queue is empty, break and finish the task
            break

        constants.downloaded = 0
        constants.queue_len = len(uids)  # number of uids in the queue
        if attempt == 1:
            # Initialize the download task
            has_task = True
        elif attempt >= max_attempt:
            failed = True
            uids_unknown += uids  # add all uids in the queue to the unknown list
            break

        try:
            with console.status(
                _get_download_banner(
                    "Make sure your internet connection is stable while we are downloading."
                ),
                refresh_per_second=100,
            ) as status:
                constants.is_downloading = True
                thread_download = threading.Thread(target=download_avatars)
                thread_download.start()  # download while fetching avatars

                task_list = list(
                    zip(
                        uids,
                        itertools.islice(  # distribute tokens evenly among instances
                            itertools.cycle(token_list), constants.queue_len
                        ),
                    )
                )
                constants.max_workers = min(32, (os.cpu_count() or 1) + 2)
                constants.delay = 0.28 / len(token_list)

                with ThreadPoolExecutor(max_workers=constants.max_workers) as pool:
                    for task_batch in np.array_split(
                        task_list, len(task_list) // (80 * len(token_list)) + 1
                    ):
                        futures = []
                        for uid, api_token in task_batch:
                            futures.append(
                                pool.submit(
                                    fetch_avatar,
                                    uid,
                                    api_token,
                                    cfg.avatar_resolution,
                                    status,
                                )
                            )

                        # Wait and finish before moving on to the next batch
                        for future in futures:
                            future.result()

                try:
                    future.result()  # type: ignore
                except AttributeError:
                    pass

                constants.is_downloading = False
                thread_download.join()

        except (ConnectionError, TimeoutError) as e:
            if attempt >= 3:
                Error(20).throw(err_type=ErrorType.WARNING)
        except InvalidTokenError as e:
            Error(21.1).throw(*e.args)
        except DiscordAPIError as e:
            Error(22).throw(*e.args)

    if uids_unknown:
        Error(23).throw(str(uids_unknown), err_type=ErrorType.WARNING)

    if has_task and not failed:
        console.print(
            "\033[A\033[2K",
            Padding(
                "[bold yellow]Avatar download complete![bold yellow]",
                (0, constants.padding, 2, constants.padding),
            ),
            sep="",
        )
        disable_console()
    # TODO: Fix while loop


if __name__ == "__main__":
    # TODO: Check merging algorithm
    # TODO: Avatar download task progress
    version_tag = "2.1"
    console.clear()
    console.set_window_title(f"Mic Drop Results {version_tag}")
    disable_console()

    # Section A: Fix console-related issues
    freeze_support()  # multiprocessing freeze support
    signal(SIGINT, SIG_IGN)  # handle KeyboardInterrupt
    atexit.register(enable_console)
    warnings.simplefilter(action="ignore", category=UserWarning)
    sys.excepthook = print_exception_hook  # avoid exiting program on exception

    # Section B: Check for missing files
    if missing_files := [
        f
        for f in (
            "data.xlsm",
            "template.pptm",
            "settings.ini",
            "token.txt",
        )
        if not abs_dir(f).is_file()
    ]:
        Error(40).throw(
            "The following files are missing:",
            "- " + "\n- ".join(missing_files),
            f"[b]Current working directory:[/b]  {MAIN_DIR}",
        )

    # Section C: Load user configurations
    cfg = Config(str(abs_dir("settings.ini")))
    n_scols = len(cfg.sort_orders)  # number of sorting columns
    avatar_mode = cfg.avatar_mode  # is subject to change later
    scheme, scheme_alt = [
        list(map(hex_to_rgb, x)) for x in (cfg.scheme, cfg.scheme_alt)
    ]

    # Section D: Parse and test tokens
    with open(abs_dir("token.txt"), "r", encoding="utf-8") as f:
        lines = f.read().splitlines()
        token_list = [line.replace('"', "").strip() for line in lines if len(line) > 70]

    if avatar_mode and not token_list:
        Error(21).throw()

    # Section E: Check for updates
    status = None
    if cfg.update_check:
        with contextlib.suppress(
            requests.exceptions.ConnectionError,
            requests.exceptions.ReadTimeout,
            KeyError,
        ):
            # Fetch the latest version and the summary of the update
            latest_tag, summary = fetch_latest_version()

            latest, current = parse_version(latest_tag, version_tag)

            if latest > current:
                status = ProgramStatus.UPDATE_AVAILABLE

                console.print(
                    Padding(
                        f"[bold yellow]Update available: Version {latest_tag}[/bold yellow]\n"
                        f"{summary}\n"
                        f"Visit release: {LATEST_RELEASE_URL}",
                        (2, constants.padding, 2, constants.padding),
                    )
                )

                time.sleep(3)
                webbrowser.open(LATEST_RELEASE_URL, new=2)
                time.sleep(3)

            elif latest < current:
                status = ProgramStatus.BETA
            else:
                status = ProgramStatus.UP_TO_DATE

    # Print a header containing information about the program

    # Normal:       Mic Drop Results (v3.10)
    #               https://github.com/SicariusBlack/mic-drop-results

    # With update:  Update v3.11
    #               A summary of the update will appear in this line.
    #               https://github.com/SicariusBlack/mic-drop-results/releases/latest/
    #
    #               Mic Drop Results (v3.10) [update available]

    status_msg = f" [{status.value}]" if status else ""
    console.print(f"[bold]Mic Drop Results{status_msg}[/bold]", justify="center")
    console.print(f"Version {version_tag}", justify="center")
    console.print(REPO_URL, justify="center")

    # Section F: Read and process the data file
    xls = pd.ExcelFile(abs_dir("data.xlsm"))
    workbook: dict[int | str, pd.DataFrame] = pd.read_excel(xls, sheet_name=None)
    xls.close()

    sheet_names = [
        match_forbidden_char.sub("", str(name)).strip()  # forbidden file name chars
        for name in workbook
    ]
    workbook = {name: list(workbook.values())[i] for i, name in enumerate(sheet_names)}

    db_prefix = "("  # signifies database tables
    database: dict[str, pd.DataFrame] = {}
    for sheet in sheet_names:
        if not sheet.startswith(db_prefix):
            continue  # exclude non-database sheets

        table = workbook[sheet]
        if table.empty or table.shape < (1, 2):  # (1 row, 2 cols) min
            continue
        database[sheet] = table

    groups: dict[str, pd.DataFrame] = {}
    for sheet in sheet_names:
        if sheet.startswith(db_prefix):
            continue  # exclude database tables

        df = workbook[sheet]
        if df.empty or df.shape < (1, n_scols):  # (1 row, n_scols cols) min
            continue

        scols = df.columns.tolist()[:n_scols]  # get sorting cols
        SHEET_INFO = (
            f"[b]Sheet name:[/b]  {sheet}\n\n"
            "See the following row(s) in data.xlsm to find out what caused the problem:"
        )

        # Exclude sheets with non-numeric sorting cols
        if any(df.loc[:, scol].dtype.kind not in "biufc" for scol in scols):
            # Get list of non-numeric vals
            str_vals = (
                df.loc[:, scols][~df.loc[:, scols].applymap(np.isreal)]
                .melt(value_name="__value")
                .dropna()["__value"]
                .tolist()
            )

            Error(60).throw(
                SHEET_INFO,
                preview_df(
                    df,
                    ~df.loc[:, scols].applymap(np.isreal).all(1),
                    n_cols=n_scols,
                    words_to_highlight=str_vals,
                ),
                err_type=ErrorType.ERROR,
            )

        # Fill nan vals within the sorting cols
        if df.loc[:, scols].isnull().values.any():
            Error(61).throw(
                SHEET_INFO,
                preview_df(
                    df,
                    df.loc[:, scols].isnull().any(axis=1),
                    n_cols=n_scols,
                    words_to_highlight=[None],
                ),
                err_type=ErrorType.WARNING,
            )

            df.loc[:, scols] = df.loc[:, scols].fillna(0)

        # Rank the slides
        df["__r"] = (
            pd.DataFrame(
                df.loc[:, scols]  # select sorting cols
                * (np.array(cfg.sort_orders) * 2 - 1)
            )  # map bool 0/1 to -1/1
            .apply(tuple, axis=1)  # type: ignore
            .rank(method="min", ascending=False)
            .astype(int)
        )

        # Sort the slides by rank
        df = df.sort_values(by="__r", ascending=True)

        # Remove .0 from whole nums
        format_int = lambda x: str(int(x)) if x % 1 == 0 else str(x)
        df.loc[:, df.dtypes == float] = df.loc[:, df.dtypes == float].applymap(
            format_int
        )

        # Replace {__sheet} with sheet name
        df["__sheet"] = sheet

        # Merge contestant database
        if database:
            process_str = lambda series: (
                series.apply(clean_name) if (series.dtype.kind == "O") else series
            )

            for table in database.values():
                df_cols = df.columns.tolist()
                db_cols = table.columns.tolist()

                anchor_col = db_cols[0]
                overlapped_cols = [
                    col for col in db_cols if (col in df_cols) and (col != anchor_col)
                ]

                if anchor_col not in df_cols:  # TODO: add a warning
                    continue

                # Copy processed vals of anchor col to '__merge_anchor'
                df["__merge_anchor"] = process_str(df[anchor_col])
                table["__merge_anchor"] = process_str(table[anchor_col])
                table = table.drop(columns=anchor_col)
                table = table.drop_duplicates("__merge_anchor")

                # Merge
                df = df.merge(table, on="__merge_anchor", how="left")
                # Note: merging wtih an existing column will produce duplicates

                for col in overlapped_cols:
                    df[col] = df[f"{col}_y"].fillna(df[f"{col}_x"])
                    df = df.drop(columns=[f"{col}_x", f"{col}_y"])

                df = df.drop(columns="__merge_anchor")

        if "__uid" not in df.columns:
            avatar_mode = False

        df["__template"] = df["__template"].fillna(1)
        df["__uid"] = df["__uid"].str.replace("_", "").str.strip()
        groups[sheet] = df

    if not groups:
        Error(68).throw()

    # Section G: Generate PowerPoint slides
    run(
        "TASKKILL /F /IM powerpnt.exe",  # kill all PowerPoint instances
        stdout=DEVNULL,
        stderr=DEVNULL,
    )

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(AVATAR_DIR, exist_ok=True)
    os.makedirs(TEMP_DIR, exist_ok=True)
    check_call(["attrib", "+H", TEMP_DIR])  # hide temp folder

    # Create Module1.bas
    with open(abs_dir(TEMP_DIR, "Module1.bas"), "w") as f:
        f.write(module1_bas)

    ### Main body
    console.print(
        Padding(
            "[bold yellow]Generating slides...[/bold yellow]\n"
            "To prevent any errors or interruptions, please avoid clicking on any PowerPoint window that pops up during the process.",
            (2, constants.padding, 2, constants.padding),
        )
    )

    thread_avatar = threading.Thread(target=_import_avatars)
    if avatar_mode:
        last_clear_dir = abs_dir(TEMP_DIR, "last_clear_avatar_cache.txt")

        try:
            with open(last_clear_dir, "r") as f:
                last_clear_time = int(f.readline())
        except (FileNotFoundError, ValueError):
            last_clear_time = 0

        if time.time() - last_clear_time > 3600 * 12:  # clear every 12 hours
            for avatar_dir in os.scandir(AVATAR_DIR):
                os.unlink(avatar_dir)

            with open(last_clear_dir, "w") as f:
                f.write(str(int(time.time())))  # update last clear time

        # Download avatars while generating slides
        thread_avatar.start()

    for sheet, df in groups.items():
        # Open template.pptm
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Presentations.Open(abs_dir("template.pptm"))

        # Minimize the window
        try:
            ppt.ActiveWindow.WindowState = 2
        except:  # catch all "no opened window" errors
            pass

        # Import macros
        try:
            ppt.VBE.ActiveVBProject.VBComponents.Import(
                abs_dir(TEMP_DIR, "Module1.bas")
            )
        except com_error as e:  # trust access not yet enabled
            if e.hresult == -2147352567:  # type: ignore
                Error(41).throw()
            else:
                raise e

        # Duplicate slides
        slides_count = ppt.Run("Count")

        # Check for invalid template IDs
        if unknown_templates := [
            x
            for x in df["__template"]
            if as_type(int, x) not in range(1, slides_count + 1)
        ]:
            showcase_cols = ["__r", "__template"]
            df_showcase = df[
                showcase_cols + [col for col in df.columns if col not in showcase_cols]
            ]
            df_showcase = df_showcase.drop_duplicates("__template").reset_index(
                drop=True
            )
            Error(71).throw(
                preview_df(
                    df_showcase,
                    n_cols=2,
                    n_cols_ext=0,
                    words_to_highlight=unknown_templates,
                )
            )

        # Duplicate initial template slides
        for template in df["__template"]:
            ppt.Run("Duplicate", template)
        ppt.Run(
            "DelSlide", *range(1, slides_count + 1)
        )  # delete initial template slides

        # Save as .pptx
        output_dir = abs_dir(OUTPUT_DIR, f"{sheet}.pptx")
        ppt.Run("SaveAs", str(output_dir))
        ppt.Quit()

        # Wait for avatars
        if avatar_mode:
            thread_avatar.join()

        # Open .pptx file and fill slides with judging data
        prs = Presentation(str(output_dir))
        for i, slide in enumerate(prs.slides):
            fill_slide(
                slide,
                {
                    k.lstrip("__"): str(
                        v
                    )  # treat program-domain vars like normal vars when replacing
                    for k, v in df.iloc[i].to_dict().items()
                },
            )

        # Save .pptx file
        prs.save(output_dir)

    # Section H: Launch the file
    inp(
        Padding(
            f"[bold yellow]Exported to {OUTPUT_DIR}[/bold yellow]\n"
            "Press Enter to open the output folder...",
            (0, constants.padding, 2, constants.padding),
        ),
        hide_text=True,
    )
    os.startfile(OUTPUT_DIR)
