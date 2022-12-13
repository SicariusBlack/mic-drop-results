import configparser
import contextlib
from ctypes import windll
from io import BytesIO
import itertools
from json import load, dump
from multiprocessing import Pool, freeze_support
import multiprocessing.popen_spawn_win32 as forking
import numpy as np
import os
from PIL import Image
import re
import requests
from signal import signal, SIGINT, SIG_IGN
from subprocess import run, DEVNULL
import sys
import time
from typing import Any
from traceback import print_exception
from urllib.request import Request, urlopen
import webbrowser

import cursor
from colorama import init, Fore, Back

import cv2
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE  # type: ignore
from pptx.slide import Slide
from pptx.util import Inches

from pywintypes import com_error
import win32com.client


class _Popen(forking.Popen):
    """Makes multiprocessing compatible with pyinstaller.

    Source:
    https://github.com/pyinstaller/pyinstaller/wiki/Recipe-Multiprocessing
    """

    def __init__(self, *args, **kw):
        if hasattr(sys, 'frozen'):
            os.putenv('_MEIPASS2', sys._MEIPASS)
        try:
            super(_Popen, self).__init__(*args, **kw)
        finally:
            if hasattr(sys, 'frozen'):
                if hasattr(os, 'unsetenv'):
                    os.unsetenv('_MEIPASS2')
                else:
                    os.putenv('_MEIPASS2', '')


forking.Popen = _Popen


class ProgressBar:
    """Creates and prints a progress bar.

    Attributes:
        progress: number of work done. Updates via the add() method.
        total: number of work to perform.
        title: title shown to the left of the progress bar.
        max_title_length: length of the longest title to ensure left
            alignment of the progress bars when there are more than
            one bar.
        bar_length: length of the progress bar in characters.
        desc: description of the task in progress shown below the
            progress bar. Updates via the set_description() method.
    """

    def __init__(self, total: int, title: str, max_title_length: int,
                 bar_length: int = 40) -> None:
        self.progress: int = 0
        self.total = total
        self.title = title
        self.max_title_length = max_title_length
        self.bar_length = bar_length
        self.desc: str = ''

    def refresh(self) -> None:
        """Reprints the progress bar with updated parameters."""
        filled_length = round(self.bar_length * self.progress / self.total)

        percents = round(100 * self.progress / self.total, 1)
        bar = '█' * filled_length + ' ' * (self.bar_length - filled_length)

        if self.progress > 0:
            sys.stdout.write('\033[2K\033[A\r')  # Delete line, move cursor up,
                                                 # and to beginning of the line
            sys.stdout.flush()

        title_right_padding = self.max_title_length - len(self.title) + 1
        sys.stdout.write(f'{self.title}{" " * title_right_padding}'
                         f'|{bar}| {self.progress}/{self.total} [{percents}%]'
                         f'{self.desc}')


        # Preview:      Merge   |████████████████████████| 7/7 [100%]
        #               Group 1 |███████████████         | 5/8 [63%]
        #               Filling in judging data


        if self.progress == self.total:
            sys.stdout.write('\033[2K\r')        # Delete line and move cursor
                                                 # to beginning of line

        sys.stdout.flush()
        
    def set_description(self, desc: str = '') -> None:
        """Sets the description shown below the progress bar."""
        self.desc = '\n' + desc
        self.refresh()

    def add(self, increment: int = 1) -> None:
        """Updates the progress by a specified increment."""
        self.progress += increment
        self.progress = min(self.progress, self.total)
        self.refresh()


def is_number(a: Any) -> bool:
    """Checks if value is a number."""
    try:
        float(a)
        return True
    except ValueError:
        return False


def as_int(a: Any) -> int | Any:
    """If possible, returns value as integer, otherwise returns value as is."""
    try:
        return int(a)
    except ValueError:
        return a


def console_style(style: str = Fore.RESET + Back.RESET + '\033[1m') -> None:
    """Sets the color and style in which the next line is printed.
    
    Args:
        color (optional): Defaults to resetting all formatting.

    Examples:
        >>> console_style(Fore.RED)
        >>> console_style(Back.YELLOW)

        To reset the color to default:

        >>> console_style()
    """
    print(style, end='')


class ErrorType:
    """Contains the string constants of error types for throw_error() function.
    
    Consts:
        E: 'ERROR', W: 'WARNING', I: 'INFO'
    """
    E = 'ERROR'
    W = 'WARNING'
    I = 'INFO'


def throw_error(*paragraphs: str, err_type: ErrorType = ErrorType.E) -> None:
    """Handles and reprints an error with additional guides and details.
    
    Prints an error message with paragraphs of extra details separated
    by single blank lines (double-spaced between). The first paragraph
    will be shown beside the err_type and will inherit the color red
    if it is an error, otherwise, in case of a warning for example,
    would be printed in yellow.

    Args:
        *paragraphs (str): 
    """
    if paragraphs:
        if err_type == ErrorType.E:
            console_style(Fore.RED)
            console_style(Back.YELLOW)
        elif err_type == ErrorType.W:
            console_style(Fore.YELLOW)

        print(f'\n\n{err_type}: {paragraphs[0]}')
        console_style()

    if len(paragraphs) > 1:
        print()
        print(*paragraphs[1:], sep='\n\n')

    if err_type == ErrorType.E:
        input_('\nPress Enter to exit the program...')
        sys.exit(1)
    else:
        input_('\nPress Enter to continue...')


def print_exception_and_exit(exc_type, exc_value, tb) -> None:
    print_exception(exc_type, exc_value, tb)
    throw_error()


def hex_to_rgb(hex_val: str) -> tuple[int, int, int]:
    return tuple(int(hex_val.lstrip('#')[i : i+2], 16) for i in (0, 2, 4))


def input_(*args, **kwargs):
    # Enable QuickEdit, thus allowing the user to copy the error message
    kernel32 = windll.kernel32
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x40|0x100))
    cursor.show()

    print(*args, **kwargs, end='')
    i = input()

    # Disable QuickEdit and Insert mode
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x00|0x100))
    cursor.hide()

    return i


def replace_text(slide: Slide, df, i, avatar_mode) -> Slide:
    """Replaces and formats text."""
    cols = df.columns.values.tolist() + ['p']

    for shape in slide.shapes:  # type: ignore
        if not shape.has_text_frame or '{' not in shape.text:
            continue

        text_frame = shape.text_frame

        for run in itertools.chain.from_iterable(
            [p.runs for p in text_frame.paragraphs]):

            for search_str in (set(re.findall(r'(?<={)(.*?)(?=})', run.text))
                               .intersection(cols)):

                # Profile picture
                if search_str == 'p':
                    # Test cases
                    if 'uid' not in cols or not avatar_mode:
                        continue

                    if pd.isnull(df['uid'].iloc[i]):
                        run.text = ''
                        continue

                    # Extract effect index and remove {p}
                    effect = as_int(run.text.strip()[3:])
                    run.text = ''

                    uid = str(df['uid'].iloc[i]).strip().replace('_', '')

                    og_path = avatar_path + '_' + uid + '.png'
                    img_path = avatar_path + str(effect) + '_' + uid + '.png'

                    if not os.path.isfile(og_path):
                        continue

                    if is_number(effect):
                        img = cv2.imread(og_path)
                        match effect:
                            case 1:
                                img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

                        cv2.imwrite(img_path, img)

                    new_shape = slide.shapes.add_picture(  # type: ignore
                        img_path, shape.left, shape.top,
                        shape.width, shape.height
                    )

                    new_shape.auto_shape_type = MSO_SHAPE.OVAL
                    old = shape._element
                    new = new_shape._element
                    old.addnext(new)
                    old.getparent().remove(old)
                    continue

                # Actual text
                repl = str(df[search_str].iloc[i])
                repl = repl if repl != 'nan' else ''  # Replace nan with empty

                run_text = run.text

                if search_str.startswith(starts):
                    run.text = repl
                else:
                    run.text = run.text.replace('{' + search_str + '}', repl)

                # Replace image links
                pattern = r'\<\<(.*?)\>\>'  # Look for <<image_links.{ext}>>
                img_link = re.findall(pattern, run.text)

                if len(img_link) > 0:
                    try:
                        img = BytesIO(requests.get(img_link[0]).content)
                        pil = Image.open(img)

                        im_width = shape.height / pil.height * pil.width
                        new_shape = slide.shapes.add_picture(  # type: ignore
                            img, shape.left + (shape.width - im_width)/2,
                            shape.top, im_width, shape.height
                        )

                        old = shape._element.addnext(new_shape._element)

                        run.text = re.sub(pattern, '', run.text)
                        text_frame.margin_left = Inches(5.2)
                    except Exception:
                        throw_error(
                            'Could not load the following image '
                           f'(Slide {i + 1}, {df["sheet"].iloc[0]}).\n{img_link[0]}',
                            'Please check your internet connection and verify that '
                            'the link leads to an image file. '
                            'It should end with an image extension like .png in most cases.',
                            err_type='WARNING')

                # Color formatting
                if not search_str.startswith(starts):
                    continue

                # Check RGB
                if (run.font.color.type == MSO_COLOR_TYPE.RGB and  # type: ignore
                    run.font.color.rgb not in [
                        RGBColor(0, 0, 0), RGBColor(255, 255, 255)]):
                    continue

                for ind, val in enumerate(range_list):
                    if is_number(repl) and float(repl) >= val:
                        if run_text.endswith('1'):
                            run.font.color.rgb = RGBColor(*color_list_light[ind])
                        else:
                            run.font.color.rgb = RGBColor(*color_list[ind])
                        break
    return slide


def get_avatar(id, api_token):
    header = {
        'Authorization': 'Bot ' + api_token
    }

    if not is_number(id):
        return None

    link = None

    try:
        response = requests.get(
            f'https://discord.com/api/v9/users/{id}', headers=header)
        link = f'https://cdn.discordapp.com/avatars/{id}/{response.json()["avatar"]}'
    except KeyError:
        if response.json()['message'] == '401: Unauthorized':
            throw_error('Invalid token. Please provide a new token in token.txt or '
                        'turn off avatar_mode in config.cfg.', response.json())

        elif response.json()['message'] == 'You are being rate-limited by the API.':
            time.sleep(response.json()['retry_after'])
            get_avatar(id, api_token)

        else:
            throw_error(response.json(), err_type='WARNING')

    except requests.exceptions.ConnectionError:
        throw_error('Unable to connect with the Discord API. Please check your '
                    'internet connection and try again.', err_type='WARNING')
    return link


def download_avatar(uid, avatar_path, api_token):
    uid = uid.strip().replace('_', '')
    img_path = avatar_path + '_' + uid.strip() + '.png'

    # Load image from link
    avatar_url = get_avatar(uid, api_token)

    if not avatar_url:
        return None

    avatar_url += '.png'

    req = urlopen(Request(avatar_url, headers={'User-Agent': 'Mozilla/5.0'}))
    arr = np.asarray(bytearray(req.read()), dtype=np.uint8)
    img = cv2.imdecode(arr, -1)

    cv2.imwrite(img_path, img)


if __name__ == '__main__':
    # Section A: Fix Command Prompt issues
    freeze_support()          # Multiprocessing freeze support
    signal(SIGINT, SIG_IGN)   # Handle KeyboardInterrupt

    # Disable QuickEdit and Insert mode
    kernel32 = windll.kernel32
    kernel32.SetConsoleMode(kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x00|0x100))

    sys.excepthook = print_exception_and_exit   # Avoid exiting the program when an error is thrown
    init()                                      # Enable ANSI escape sequences
    cursor.hide()                               # Hide cursor


    # Section B: Check for missing files
    if missing := [f for f in (

            'config.cfg', 'data.xlsx', 'template.pptm', 'Module1.bas', 'token.txt'

        ) if not os.path.isfile(f)]:
        throw_error('The following files are missing. Please review the documentation for more '
                    'information regarding file requirements.', '\n'.join(missing))


    # Section C: Load config.cfg
    config = load(open('config.json'))

    # Store config variables as local variables
    range_list = config['format']['ranges'][::-1]
    color_list = config['format']['colors'][::-1]
    color_list_light = config['format']['colors_light'][::-1]
    starts = config['format']['starts_with']
    avatar_mode = config['avatars']
    last_clear = config['last_clear_avatar_cache']

    with open('token.txt') as f:
        token_list = f.read().splitlines()
        token_list = [i.strip() for i in token_list if len(i) > 62]

    if not token_list and avatar_mode:
        throw_error('Please provide a valid bot token in token.txt or turn off avatar mode in config.json.')

    color_list = list(map(hex_to_rgb, color_list))
    color_list_light = list(map(hex_to_rgb, color_list_light))


    # Section D: Check for updates
    status, url = '', ''

    with contextlib.suppress(requests.exceptions.ConnectionError, KeyError):
        if config['update_check']:
            response = requests.get('https://api.github.com/repos/banz04/mic-drop-results/releases/latest', timeout=3)

            raw_ver = response.json()['tag_name'][1:]
            version, config_ver = [tuple(map(int, v.split('.'))) for v in (
                raw_ver, config['version']
            )]

            if version > config_ver:
                console_style(Fore.YELLOW)
                print(f'Update v{raw_ver}')
                print(response.json()['body'].partition('\n')[0])

                url = 'https://github.com/banz04/mic-drop-results/releases/latest/'
                print(url + '\n')
                webbrowser.open(url, new=2)
                console_style()

                status = 'update available'
            elif version < config_ver:
                status = 'beta'
            else:
                status = 'latest'

            status = f' [{status}]'

    # Print a header containing information about the program

    # Preview:      Mic Drop Results (vX.1) [latest]
    #               https://github.com/banz04/mic-drop-results


    # With update:  Update vX.1
    #               A summary of the update will appear in this line.
    #               https://github.com/banz04/mic-drop-results/releases/latest/
    #
    #               Mic Drop Results (vX.0) [update available]


    print(f'Mic Drop Results (v{config["version"]}){status}')
    console_style()

    if 'update available' not in status:
        url = 'https://github.com/banz04/mic-drop-results'
        print(url)


    # Section E: Process the data
    folder_path = os.getcwd() + '\\'
    output_path = folder_path + 'output\\'
    avatar_path = folder_path + 'avatars\\'

    xls = pd.ExcelFile('data.xlsx')

    sheetnames_raw = xls.sheet_names
    sheetnames = [re.sub(r'[\\\/:"*?<>|]+',  # Forbidden file name characters
        '', sheet) for sheet in sheetnames_raw]

    data = {}

    db_list = []
    for sheet in sheetnames_raw:
        if sheet.startswith('_'):
            df = pd.read_excel(xls, sheet)

            # Validate shape
            if df.empty or df.shape < (1, 2):
                continue

            db_list.append(df)

    SHARING_VIOLATION = '\033[33mNOTE: Please exit the program before modifying data.xlsx or ' \
                        'Microsoft Excel will throw a Sharing Violation error.\033[39m'

    for i, sheet in enumerate(sheetnames_raw):
        df = pd.read_excel(xls, sheet)

        # Validate shape
        if df.empty or df.shape < (1, 2):
            continue

        # Exclude database sheets
        if sheet.startswith('_'):
            continue

        # Exclude sheets with first two columns where data is not numeric
        if sum(df.iloc[:, i].dtype.kind in 'biufc' for i in range(2)) < 2:
            throw_error(f'Invalid data type. The following rows of {sheet} contain strings '
                         'instead of the supposed numeric data type within the first two columns. '
                         'The sheet will be excluded if you proceed on.',

                df[~df.iloc[:, :2].applymap(np.isreal).all(1)],

                err_type='warning'
            )

            continue

        # Replace NaN values within the first two columns with 0
        if df.iloc[:, :2].isnull().values.any():
            throw_error(f'The following rows of {sheet} contain empty values '
                'within the first two columns.',

                df[df.iloc[:, :2].isnull().any(axis=1)],

                'You may exit this program and modify your data or proceed on with '
                'these empty values substituted with 0.', SHARING_VIOLATION,

                err_type='warning'
            )

            df.iloc[:, :2] = df.iloc[:, :2].fillna(0)

        # Check for cases where avg and std are the same (hold the same rank)
        df['r'] = pd.DataFrame(zip(df.iloc[:, 0], df.iloc[:, 1] * -1)) \
                    .apply(tuple, axis=1).rank(method='min', ascending=False).astype(int)

        # Sort the slides
        df = df.sort_values(by='r', ascending=True)

        # Remove .0 from whole numbers
        format_number = lambda x: str(int(x)) if x % 1 == 0 else str(x)
        df.loc[:, df.dtypes == float] = df.loc[:, df.dtypes == float].applymap(format_number)

        # Replace {sheet} with sheet name
        df['sheet'] = sheet

        # Merge contestant database
        clean_name = lambda x: x.str.lower().str.strip() if(x.dtype.kind == 'O') else x
        if db_list:
            for db in db_list:
                df_cols = df.columns.values.tolist()
                db_cols = db.columns.values.tolist()
                merge_col = db_cols[0]

                # Use merge for non-existing columns
                df = df.merge(db[[merge_col] + [i for i in db_cols if i not in df_cols]],
                    left_on=clean_name(df[merge_col]), right_on=clean_name(db[merge_col]), how='left')

                df.loc[:, merge_col] = df[merge_col + '_x']
                df.drop(['key_0', merge_col + '_x', merge_col + '_y'], axis=1, inplace=True)

                # Use update for existing columns
                df['update_index'] = clean_name(df[merge_col])
                df = df.set_index('update_index')

                db['update_index'] = clean_name(db[merge_col])
                db = db.set_index('update_index')
                db_cols.remove(merge_col)

                update_cols = [i for i in db_cols if i in df_cols]
                df.update(db[update_cols])
                df.reset_index(drop=True, inplace=True)

        if 'uid' not in df.columns.values.tolist(): avatar_mode = 0

        # Fill in missing templates
        df['template'].fillna(1, inplace=True)

        data[sheetnames[i]] = df

    if not data:
        throw_error(f'No valid sheet found in {folder_path}data.xlsx')


    # Section F: Generate PowerPoint slides
    print('\nGenerating slides...')
    print('Please do not click on any PowerPoint windows that may show up in the process.\n')

    # Kill all PowerPoint instances
    run('TASKKILL /F /IM powerpnt.exe', stdout=DEVNULL, stderr=DEVNULL)

    # Open template presentation
    os.makedirs(output_path, exist_ok=True)
    os.makedirs(avatar_path, exist_ok=True)

    # Clear cache
    if time.time() - last_clear > 1800:  # Clears every hour
        for f in os.scandir(avatar_path):
            os.unlink(f)

        # Update last clear time
        with open('config.json', 'w') as f:
            config['last_clear_avatar_cache'] = int(time.time())
            dump(config, f, indent=4)

    # Download avatars with parallel processing
    attempt = 0
    pool = Pool(3)

    while avatar_mode:
        uid_list = []

        for df in data.values():
            if df['uid'].dtype.kind in 'biufc':
                throw_error('The \'uid\' column has a numeric data type instead of the supposed string data type.',
                            'Please exit the program and add an underscore before every user ID.', SHARING_VIOLATION)

            uid_list += [id for id in df['uid'] if not pd.isnull(id) and not os.path.isfile(avatar_path + id.strip() + '.png')]

        if len(uid_list) == 0:
            break

        if attempt > 0 and attempt <= 3:
            print(f'Unable to download the profile pictures of the following users. Retrying {attempt}/3',
                    uid_list, sep='\n', end='\n\n')
        elif attempt > 3:
            throw_error('Failed to download the profile pictures of the following users. Please verify that their user IDs are correct.',
                    str(uid_list), err_type='warning')

        pool.starmap(download_avatar, zip(uid_list,
            [avatar_path] * len(uid_list), itertools.islice(itertools.cycle(token_list), len(uid_list))))

        attempt += 1

    pool.close()
    pool.join()

    for k, df in data.items():
        bar = ProgressBar(8, title=k, max_title_length=max(map(len, data.keys())))

        # Open template presentation
        bar.set_description('Opening template.pptm')
        ppt = win32com.client.Dispatch('PowerPoint.Application')
        ppt.Presentations.Open(f'{folder_path}template.pptm')
        bar.add()

        # Import macros
        bar.set_description('Importing macros')

        try:
            ppt.VBE.ActiveVBProject.VBComponents.Import(f'{folder_path}Module1.bas')
        except com_error as e:
            if e.hresult == -2147352567:  # type: ignore
                # Warns the user about trust access error
                throw_error('Please open PowerPoint, look up Trust Center Settings, '
                            'and make sure Trust access to the VBA project object model is enabled.')
            else:
                raise e

        bar.add()

        # Duplicate slides
        bar.set_description('Duplicating slides')
        slides_count = ppt.Run('Count')

        # Duplicate slides
        for t in df.loc[:, 'template']:
            if as_int(t) not in range(1, slides_count + 1):
                throw_error(f'Template {t} does not exist (error originated from the following sheet: {k}).',
                            f'Please exit the program and modify the \'template\' column of {k}.', SHARING_VIOLATION)

            ppt.Run('Duplicate', t)

        bar.add()

        # Delete template slides when done
        ppt.Run('DelSlide', *range(1, slides_count + 1))
        bar.add()

        # Save as output file
        bar.set_description('Saving templates')
        output_filename = f'{k}.pptx'

        ppt.Run('SaveAs', f'{output_path}{output_filename}')
        bar.add()

        run('TASKKILL /F /IM powerpnt.exe', stdout=DEVNULL, stderr=DEVNULL)
        bar.add()

        # Replace text
        bar.set_description('Filling in judging data')
        prs = Presentation(output_path + output_filename)

        for i, slide in enumerate(prs.slides):
            replace_text(slide, df, i, avatar_mode)
        bar.add()

        # Save
        bar.set_description(f'Saving as {output_path + output_filename}')
        prs.save(output_path + output_filename)
        bar.add()


    # Section G: Launch the file
    print(f'\nExported to {output_path}')

    # Enable QuickEdit
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x40|0x100))

    input_('Press Enter to open the output folder...')
    os.startfile(output_path)
